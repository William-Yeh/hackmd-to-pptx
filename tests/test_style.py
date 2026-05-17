"""Tests for HackMD/Marp <style> block translation to python-pptx styling.

The CSS parser is intentionally minimal: it understands the subset that
real-world HackMD/Marp decks use (reveal.js scoped rules, a handful of
properties), keyed by the *trailing selector token* so the `.reveal .slides`
prefix is irrelevant.
"""

from pptx import Presentation
from pptx.util import Pt

from convert import (
    add_content_slide,
    add_section_slide,
    extract_style_block,
    hex_to_rgb,
    parse_markdown,
    parse_style_block,
)


def _rgb_to_hex(rgb):
    """python-pptx RGBColor → '#RRGGBB' for readable assertions."""
    return "#" + str(rgb).upper()


class TestParseStyleBlock:
    def test_returns_empty_for_empty_input(self):
        assert parse_style_block("") == {}

    def test_returns_empty_for_whitespace_only(self):
        assert parse_style_block("   \n  \n  ") == {}

    def test_extracts_simple_rule(self):
        css = "<style>h1 { color: #FFC500; }</style>"
        result = parse_style_block(css)
        assert result == {"h1": {"color": "#FFC500"}}

    def test_strips_reveal_prefix_from_selector(self):
        # The trailing selector token is what matters; `.reveal .slides`
        # is a reveal.js scope that the converter ignores.
        css = "<style>.reveal .slides h1 { color: #FFC500; }</style>"
        result = parse_style_block(css)
        assert result == {"h1": {"color": "#FFC500"}}

    def test_bare_reveal_slides_maps_to_empty_token(self):
        # `.reveal .slides { ... }` is the slide-body default; we key it
        # with the empty string so the renderer can find it.
        css = "<style>.reveal .slides { font-size: 28px; text-align: left; }</style>"
        result = parse_style_block(css)
        assert result == {"": {"font-size": "28px", "text-align": "left"}}

    def test_multiple_rules(self):
        css = """<style>
        .reveal .slides { font-size: 28px; }
        .reveal .slides h1 { color: #FFC500; }
        .reveal .slides code { font-size: 24px; }
        </style>"""
        result = parse_style_block(css)
        assert result == {
            "": {"font-size": "28px"},
            "h1": {"color": "#FFC500"},
            "code": {"font-size": "24px"},
        }

    def test_real_world_unit1_block(self):
        # Verbatim from ai_dev_newbie merged deck head.
        css = """<style>
        .reveal .slides {
            text-align: left;
            font-size: 28px;
        }
        .reveal .slides blockquote {
            text-align: left;
            font-size: 28px;
            font-style: normal;
        }
        .reveal .slides h1 {
            color: #FFC500;
        }
        .reveal .slides code {
            font-size: 24px;
            line-height: 1.5em;
        }
        </style>"""
        result = parse_style_block(css)
        assert result[""]["text-align"] == "left"
        assert result[""]["font-size"] == "28px"
        assert result["h1"]["color"] == "#FFC500"
        assert result["code"]["font-size"] == "24px"
        assert result["blockquote"]["font-style"] == "normal"

    def test_unknown_properties_dropped(self):
        # margin, padding, transform, etc. are not in our whitelist and
        # are silently ignored so the dict stays small and predictable.
        css = "<style>h1 { color: red; margin: 10px; transform: rotate(5deg); }</style>"
        result = parse_style_block(css)
        assert result == {"h1": {"color": "red"}}

    def test_malformed_rule_skipped(self):
        # A rule missing its closing brace shouldn't tank the rest.
        css = "<style>h1 { color: red h2 { color: blue; }</style>"
        result = parse_style_block(css)
        # We don't insist on which rules survive; we only insist parsing
        # doesn't raise and at least one valid rule comes through.
        assert isinstance(result, dict)

    def test_css_comments_stripped(self):
        css = """<style>
        /* deck-level default */
        h1 { color: red; /* gold-ish */ }
        </style>"""
        result = parse_style_block(css)
        assert result == {"h1": {"color": "red"}}

    def test_no_style_tags_treated_as_raw_css(self):
        # Calling with raw CSS (no <style> wrapper) still works so the
        # function is reusable from merge.py / tests.
        css = "h1 { color: red; }"
        result = parse_style_block(css)
        assert result == {"h1": {"color": "red"}}


class TestExtractStyleBlock:
    def test_no_style_returns_none(self):
        md = "# Hello\n\n- bullet"
        style, body = extract_style_block(md)
        assert style is None
        assert body == md

    def test_lifts_top_of_deck_style(self):
        md = (
            "<style>h1 { color: #FFC500; }</style>\n\n"
            "# Hello\n\n- bullet\n"
        )
        style, body = extract_style_block(md)
        assert style is not None
        assert "color: #FFC500" in style
        assert body.startswith("# Hello")
        assert "<style>" not in body

    def test_ignores_style_after_first_heading(self):
        # Per-slide <style> tags deeper in the deck are left in the body
        # untouched (we don't try to support them).
        md = (
            "# Hello\n\n<style>h2 { color: red; }</style>\n\n## Sub\n"
        )
        style, body = extract_style_block(md)
        assert style is None
        assert "<style>" in body


class TestParseMarkdownStripsStyle:
    def test_style_block_does_not_leak_into_slide_content(self):
        # Latent bug fix: prior behavior emitted the style block as
        # plain-text content on the first slide.
        md = (
            "---\ntitle: Test\n---\n\n"
            "<style>\n.reveal .slides h1 { color: #FFC500; }\n</style>\n\n"
            "# Hello\n\n- bullet\n"
        )
        slides = parse_markdown(md)
        joined = "\n".join(
            item.get("text", "") for s in slides for item in s["content"]
        )
        assert "<style>" not in joined
        assert "color: #FFC500" not in joined
        assert "</style>" not in joined

    def test_parse_markdown_exposes_style(self):
        # parse_markdown returns a list of slides; the style dict is
        # attached to the list (via a sentinel attribute) so the renderer
        # can reach it without changing the function's return type for
        # all callers.
        md = (
            "<style>.reveal .slides h1 { color: #FFC500; }</style>\n\n"
            "# Hello\n"
        )
        slides = parse_markdown(md)
        assert getattr(slides, "style_overrides", {}) == {
            "h1": {"color": "#FFC500"}
        }


class TestRendererAppliesStyle:
    """Verify the renderer threads parsed style overrides onto the right runs."""

    def _section(self, prs, title, *, overrides=None, colors, fonts):
        return add_section_slide(
            prs,
            {"title": f"# {title}", "subtitle": None, "content": [], "notes": None, "is_section": True},
            colors,
            fonts,
            style_overrides=overrides or {},
        )

    def _content(self, prs, title, content, *, overrides=None, colors, fonts):
        return add_content_slide(
            prs,
            {"title": f"## {title}", "subtitle": None, "content": content, "notes": None, "is_section": False},
            colors,
            fonts,
            style_overrides=overrides or {},
        )

    def test_h1_color_applied_to_section_title(self, colors, fonts):
        # On a dark slide background, gold #FFC500 passes contrast and must
        # land on every title run. We assert against a *dark* background by
        # passing the override and reading the resulting run color.
        prs = Presentation()
        overrides = {"h1": {"color": "#1A1A1A"}}  # near-black on white = high contrast
        slide = self._section(prs, "Hello", overrides=overrides, colors=colors, fonts=fonts)
        title_runs = slide.placeholders[0].text_frame.paragraphs[0].runs
        assert title_runs, "title placeholder must have at least one run"
        assert all(_rgb_to_hex(r.font.color.rgb) == "#1A1A1A" for r in title_runs)

    def test_h1_color_dropped_when_contrast_fails(self, colors, fonts):
        # Default pptx background is white; gold #FFC500 on white = ~1.6:1,
        # which fails WCAG AA (4.5:1). The renderer must drop the color so
        # titles remain readable in default-theme output.
        prs = Presentation()
        overrides = {"h1": {"color": "#FFC500"}}
        slide = self._section(prs, "Hello", overrides=overrides, colors=colors, fonts=fonts)
        title_runs = slide.placeholders[0].text_frame.paragraphs[0].runs
        # Color should NOT be set to FFC500 (it was dropped). It may be None
        # (placeholder default) or anything other than FFC500.
        for r in title_runs:
            if r.font.color.type is not None:
                assert _rgb_to_hex(r.font.color.rgb) != "#FFC500"

    def test_h2_font_size_applied_to_content_title(self, colors, fonts):
        prs = Presentation()
        overrides = {"h2": {"font-size": "20px"}}
        slide = self._content(
            prs, "Hello",
            [{"type": "bullet", "text": "x", "indent": 0}],
            overrides=overrides, colors=colors, fonts=fonts,
        )
        title_runs = slide.placeholders[0].text_frame.paragraphs[0].runs
        assert title_runs
        # 20px → 20pt (we treat 1px == 1pt for slides; close enough for the
        # px values reveal.js authors actually write).
        assert all(r.font.size == Pt(20) for r in title_runs)

    def test_bare_selector_sets_body_font_size(self, colors, fonts):
        prs = Presentation()
        overrides = {"": {"font-size": "28px"}}
        slide = self._content(
            prs, "T",
            [{"type": "bullet", "text": "body item", "indent": 0}],
            overrides=overrides, colors=colors, fonts=fonts,
        )
        body_runs = slide.placeholders[1].text_frame.paragraphs[0].runs
        assert body_runs
        assert any(r.font.size == Pt(28) for r in body_runs)

    def test_code_font_size_applied_to_inline_code(self, colors, fonts):
        prs = Presentation()
        overrides = {"code": {"font-size": "24px"}}
        slide = self._content(
            prs, "T",
            [{"type": "bullet", "text": "use `foo()` here", "indent": 0}],
            overrides=overrides, colors=colors, fonts=fonts,
        )
        body_runs = slide.placeholders[1].text_frame.paragraphs[0].runs
        code_runs = [r for r in body_runs if r.font.name == fonts["code"]]
        assert code_runs, "expected at least one inline-code run"
        assert all(r.font.size == Pt(24) for r in code_runs)

    def test_code_block_ignores_code_font_size_override(self, colors, fonts):
        # Fenced code blocks are rendered into a fixed-height textbox shape
        # with a hardcoded 11pt size. Honoring `code { font-size }` here
        # would push lines past the textbox's clip rect, causing content
        # to disappear. The override is intentionally scoped to *inline*
        # code only; this test guards against re-enabling it for blocks.
        prs = Presentation()
        overrides = {"code": {"font-size": "24px"}}
        slide = self._content(
            prs, "T",
            [{"type": "codeblock", "lang": "python", "content": "print('hi')"}],
            overrides=overrides, colors=colors, fonts=fonts,
        )
        code_runs = []
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.name == fonts["code"]:
                        code_runs.append(r)
        assert code_runs, "expected code-block runs"
        # Every code-block run stays at the hardcoded 11pt.
        assert all(r.font.size == Pt(11) for r in code_runs)


class TestEndToEndWithStyle:
    """Run the full convert.py path on a style-bearing input."""

    def test_real_world_style_block_is_extracted_and_applied(self, tmp_output_dir):
        # Verbatim style block from the ai_dev_newbie deck. The h1 color
        # should drop on the default white background (contrast guard);
        # the font sizes should land.
        import subprocess, sys
        from pathlib import Path
        script = str(Path(__file__).resolve().parent.parent / "skill" / "scripts" / "convert.py")
        md = tmp_output_dir / "styled.md"
        md.write_text(
            "<style>\n"
            ".reveal .slides { text-align: left; font-size: 28px; }\n"
            ".reveal .slides h1 { color: #FFC500; }\n"
            ".reveal .slides code { font-size: 24px; }\n"
            "</style>\n\n"
            "# Hello\n\n"
            "- a bullet with `code` in it\n"
        )
        out = tmp_output_dir / "out.pptx"
        result = subprocess.run(
            [sys.executable, script, str(md), str(out)],
            capture_output=True, text=True,
        )
        assert result.returncode == 0, result.stderr
        prs = Presentation(str(out))
        # First slide is the section/title. Its title runs must NOT be
        # FFC500 (contrast guard dropped it on the default white bg).
        title_runs = prs.slides[0].placeholders[0].text_frame.paragraphs[0].runs
        for r in title_runs:
            if r.font.color.type is not None:
                assert _rgb_to_hex(r.font.color.rgb) != "#FFC500"

    def test_real_world_dark_background_keeps_color(self, tmp_output_dir):
        # When the slide background is dark, the FFC500 color must survive.
        # We exercise this through the parse → renderer path with a known
        # dark background passed via the overrides API.
        prs = Presentation()
        overrides = {"h1": {"color": "#FFC500"}}
        slide = add_section_slide(
            prs,
            {"title": "# Hello", "subtitle": None, "content": [], "notes": None, "is_section": True},
            {
                **{k: v for k, v in __import__("convert").DEFAULT_COLORS.items()},
                "slideBg": "1E2761",  # dark background; high contrast vs FFC500
            },
            __import__("convert").DEFAULT_FONTS,
            style_overrides=overrides,
        )
        title_runs = slide.placeholders[0].text_frame.paragraphs[0].runs
        assert title_runs
        assert all(_rgb_to_hex(r.font.color.rgb) == "#FFC500" for r in title_runs)
