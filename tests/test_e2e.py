import json
import subprocess
import sys
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation

CONVERT_SCRIPT = str(Path(__file__).resolve().parent.parent / "skill" / "scripts" / "convert.py")
DEMO_MD = str(Path(__file__).resolve().parent.parent / "examples" / "demo.md")


class TestEndToEnd:
    def test_demo_produces_pptx(self, tmp_output_dir):
        out = str(tmp_output_dir / "output.pptx")
        result = subprocess.run(
            [sys.executable, CONVERT_SCRIPT, DEMO_MD, out],
            capture_output=True, text=True,
        )
        assert result.returncode == 0, result.stderr
        assert Path(out).exists()

    def test_output_is_valid_zip(self, tmp_output_dir):
        out = str(tmp_output_dir / "output.pptx")
        subprocess.run(
            [sys.executable, CONVERT_SCRIPT, DEMO_MD, out],
            capture_output=True, text=True,
        )
        assert zipfile.is_zipfile(out)

    def test_slide_count(self, tmp_output_dir):
        out = str(tmp_output_dir / "output.pptx")
        subprocess.run(
            [sys.executable, CONVERT_SCRIPT, DEMO_MD, out],
            capture_output=True, text=True,
        )
        prs = Presentation(out)
        assert len(prs.slides) >= 8

    def test_first_slide_title(self, tmp_output_dir):
        out = str(tmp_output_dir / "output.pptx")
        subprocess.run(
            [sys.executable, CONVERT_SCRIPT, DEMO_MD, out],
            capture_output=True, text=True,
        )
        prs = Presentation(out)
        first_slide = prs.slides[0]
        title_text = first_slide.placeholders[0].text_frame.text
        assert "My Awesome Presentation" in title_text

    def test_missing_input_exits_1(self, tmp_output_dir):
        result = subprocess.run(
            [sys.executable, CONVERT_SCRIPT, "/nonexistent/file.md"],
            capture_output=True, text=True,
        )
        assert result.returncode == 1

    def test_table_renders_as_table_shape(self, tmp_output_dir):
        md_file = tmp_output_dir / "table.md"
        md_file.write_text(
            "# Risks\n\n"
            "## Risk worksheet\n\n"
            "| 類型 | 我看到的問題 | 為什麼這是風險？ |\n"
            "|---|---|---|\n"
            "| 看不懂 |  |  |\n"
            "| 需求不明 |  |  |\n"
            "| 驗證不足 |  |  |\n"
        )
        out = str(tmp_output_dir / "table.pptx")
        result = subprocess.run(
            [sys.executable, CONVERT_SCRIPT, str(md_file), out],
            capture_output=True, text=True,
        )
        assert result.returncode == 0, result.stderr
        prs = Presentation(out)
        # Locate the slide that actually contains the table (section slide
        # has no body content; the worksheet slide is the one with shapes).
        table_slides = [s for s in prs.slides if any(sh.has_table for sh in s.shapes)]
        assert len(table_slides) == 1, "expected exactly one slide with a table shape"
        tables = [sh.table for sh in table_slides[0].shapes if sh.has_table]
        table = tables[0]
        assert len(table.rows) == 4  # header + 3 body rows
        assert len(table.columns) == 3
        assert table.cell(0, 0).text_frame.text == "類型"
        assert table.cell(1, 0).text_frame.text == "看不懂"
        # Empty body cell is preserved as empty text
        assert table.cell(1, 1).text_frame.text == ""

        # Header styling: accent fill + bold runs (locks in the design choice)
        from convert import DEFAULT_COLORS, hex_to_rgb
        header_cell = table.cell(0, 0)
        assert header_cell.fill.fore_color.rgb == hex_to_rgb(DEFAULT_COLORS["accent"])
        header_runs = header_cell.text_frame.paragraphs[0].runs
        assert header_runs, "header cell should have at least one run"
        assert all(r.font.bold for r in header_runs)

    def test_custom_config(self, tmp_output_dir):
        md_file = tmp_output_dir / "test.md"
        md_file.write_text("# Custom Title\n\n- item one\n")
        config = tmp_output_dir / "config.json"
        config.write_text(json.dumps({"colors": {"primary": "FF0000"}}))
        out = str(tmp_output_dir / "out.pptx")

        result = subprocess.run(
            [sys.executable, CONVERT_SCRIPT, str(md_file), out],
            capture_output=True, text=True,
        )
        assert result.returncode == 0
        assert Path(out).exists()
