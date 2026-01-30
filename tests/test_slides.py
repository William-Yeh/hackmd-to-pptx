from pptx import Presentation
from pptx.oxml.ns import qn

from convert import (
    add_content_slide,
    add_formatted_runs,
    add_section_slide,
    disable_bullet,
)


def _make_prs():
    prs = Presentation()
    return prs


class TestAddSectionSlide:
    def test_creates_slide(self, colors, fonts):
        prs = _make_prs()
        data = {
            "title": "# Section Title",
            "subtitle": "Sub here",
            "content": [],
            "notes": None,
            "is_section": True,
        }
        slide = add_section_slide(prs, data, colors, fonts)
        assert len(prs.slides) == 1

    def test_title_text(self, colors, fonts):
        prs = _make_prs()
        data = {
            "title": "# Hello World",
            "subtitle": None,
            "content": [],
            "notes": None,
            "is_section": True,
        }
        slide = add_section_slide(prs, data, colors, fonts)
        # Title placeholder should contain "Hello World"
        title_text = slide.placeholders[0].text_frame.text
        assert "Hello World" in title_text

    def test_notes(self, colors, fonts):
        prs = _make_prs()
        data = {
            "title": "# T",
            "subtitle": None,
            "content": [],
            "notes": "My notes",
            "is_section": True,
        }
        slide = add_section_slide(prs, data, colors, fonts)
        assert slide.notes_slide.notes_text_frame.text == "My notes"


class TestAddContentSlide:
    def test_bullet_content(self, colors, fonts):
        prs = _make_prs()
        data = {
            "title": "## Bullets",
            "subtitle": None,
            "content": [
                {"type": "bullet", "text": "item one", "indent": 0},
                {"type": "bullet", "text": "item two", "indent": 0},
            ],
            "notes": None,
            "is_section": False,
        }
        slide = add_content_slide(prs, data, colors, fonts)
        assert len(prs.slides) == 1

    def test_numbered_content(self, colors, fonts):
        prs = _make_prs()
        data = {
            "title": "## Nums",
            "subtitle": None,
            "content": [
                {"type": "numbered", "number": "1", "text": "first", "indent": 0},
            ],
            "notes": None,
            "is_section": False,
        }
        slide = add_content_slide(prs, data, colors, fonts)
        body = slide.placeholders[1].text_frame.text
        assert "1." in body

    def test_code_block_adds_shapes(self, colors, fonts):
        prs = _make_prs()
        data = {
            "title": "## Code",
            "subtitle": None,
            "content": [
                {"type": "codeblock", "lang": "python", "content": "print('hi')"},
            ],
            "notes": None,
            "is_section": False,
        }
        slide = add_content_slide(prs, data, colors, fonts)
        # Code blocks add extra shapes (rect + textbox) beyond placeholders
        assert len(slide.shapes) > 2

    def test_notes_on_content(self, colors, fonts):
        prs = _make_prs()
        data = {
            "title": "## T",
            "subtitle": None,
            "content": [{"type": "bullet", "text": "x", "indent": 0}],
            "notes": "Note text",
            "is_section": False,
        }
        slide = add_content_slide(prs, data, colors, fonts)
        assert slide.notes_slide.notes_text_frame.text == "Note text"


class TestAddFormattedRuns:
    def test_bold_run(self, colors, fonts):
        prs = _make_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        tf = slide.placeholders[1].text_frame
        p = tf.paragraphs[0]
        add_formatted_runs(p, "**bold**", colors, fonts)
        assert any(r.font.bold for r in p.runs)

    def test_italic_run(self, colors, fonts):
        prs = _make_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        tf = slide.placeholders[1].text_frame
        p = tf.paragraphs[0]
        add_formatted_runs(p, "_italic_", colors, fonts)
        assert any(r.font.italic for r in p.runs)

    def test_code_run_uses_code_font(self, colors, fonts):
        prs = _make_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        tf = slide.placeholders[1].text_frame
        p = tf.paragraphs[0]
        add_formatted_runs(p, "`code`", colors, fonts)
        assert any(r.font.name == fonts["code"] for r in p.runs)


class TestDisableBullet:
    def test_adds_bu_none(self, colors, fonts):
        prs = _make_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        tf = slide.placeholders[1].text_frame
        p = tf.paragraphs[0]
        disable_bullet(p)
        pPr = p._p.find(qn("a:pPr"))
        bu_none = pPr.find(qn("a:buNone"))
        assert bu_none is not None
