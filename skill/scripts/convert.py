#!/usr/bin/env -S uv run --script
# Inline metadata for standalone execution (uv run skill/scripts/convert.py).
# For project-based install, see pyproject.toml.
# /// script
# requires-python = ">=3.9"
# dependencies = [
#   "python-pptx>=0.6.21",
#   "lxml>=4.9.0",
#   "PyYAML>=6.0",
# ]
# ///
"""
HackMD/Marp Markdown to PowerPoint Converter
Uses python-pptx for proper slide master/layout placeholder support
Full Markdown formatting, hyperlinks, sections, and syntax highlighting
"""

__author__ = "William Yeh"
__email__ = "william.pjyeh@gmail.com"

import re
import sys
import json
import zipfile
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree

# Try to import yaml
try:
    import yaml
    HAS_YAML = True
except ImportError:
    HAS_YAML = False

# Default configuration
DEFAULT_COLORS = {
    'primary': '1E2761',
    'secondary': 'CADCFC', 
    'accent': '0891B2',
    'white': 'FFFFFF',
    'lightBg': 'F8FAFC',
    'darkText': '1E293B',
    'mutedText': '64748B',
    'codeBlock': 'F1F5F9',
    'codeBorder': 'E2E8F0',
    # Syntax highlighting colors
    'syntaxKeyword': '7C3AED',    # Purple
    'syntaxString': '059669',     # Green
    'syntaxComment': '6B7280',    # Gray
    'syntaxNumber': 'DC2626',     # Red
    'syntaxFunction': '2563EB',   # Blue
    'syntaxType': 'D97706',       # Orange
    'syntaxDiffAdd': '059669',    # Green
    'syntaxDiffDel': 'DC2626',    # Red
}

DEFAULT_FONTS = {
    'header': 'Trebuchet MS',
    'body': 'Calibri',
    'code': 'Consolas',
}

# Language keywords for syntax highlighting
SYNTAX_KEYWORDS = {
    'python': ['def', 'class', 'import', 'from', 'return', 'if', 'elif', 'else', 'for', 'while', 
               'try', 'except', 'finally', 'with', 'as', 'lambda', 'yield', 'raise', 'pass',
               'break', 'continue', 'and', 'or', 'not', 'in', 'is', 'None', 'True', 'False',
               'async', 'await', 'global', 'nonlocal', 'assert', 'del'],
    'javascript': ['function', 'const', 'let', 'var', 'return', 'if', 'else', 'for', 'while',
                   'switch', 'case', 'break', 'continue', 'try', 'catch', 'finally', 'throw',
                   'new', 'this', 'class', 'extends', 'import', 'export', 'default', 'from',
                   'async', 'await', 'yield', 'typeof', 'instanceof', 'null', 'undefined',
                   'true', 'false', 'of', 'in'],
    'java': ['public', 'private', 'protected', 'class', 'interface', 'extends', 'implements',
             'static', 'final', 'void', 'return', 'if', 'else', 'for', 'while', 'do', 'switch',
             'case', 'break', 'continue', 'try', 'catch', 'finally', 'throw', 'throws', 'new',
             'this', 'super', 'import', 'package', 'null', 'true', 'false', 'instanceof'],
    'go': ['func', 'package', 'import', 'return', 'if', 'else', 'for', 'range', 'switch', 'case',
           'default', 'break', 'continue', 'go', 'defer', 'select', 'chan', 'map', 'struct',
           'interface', 'type', 'const', 'var', 'nil', 'true', 'false', 'make', 'new', 'append'],
    'rust': ['fn', 'let', 'mut', 'const', 'pub', 'mod', 'use', 'struct', 'enum', 'impl', 'trait',
             'return', 'if', 'else', 'for', 'while', 'loop', 'match', 'break', 'continue',
             'move', 'ref', 'self', 'Self', 'true', 'false', 'None', 'Some', 'Ok', 'Err', 'async', 'await'],
    'sql': ['SELECT', 'FROM', 'WHERE', 'INSERT', 'UPDATE', 'DELETE', 'CREATE', 'DROP', 'ALTER',
            'TABLE', 'INDEX', 'JOIN', 'LEFT', 'RIGHT', 'INNER', 'OUTER', 'ON', 'AND', 'OR', 'NOT',
            'IN', 'LIKE', 'BETWEEN', 'IS', 'NULL', 'ORDER', 'BY', 'GROUP', 'HAVING', 'LIMIT',
            'OFFSET', 'UNION', 'AS', 'DISTINCT', 'COUNT', 'SUM', 'AVG', 'MAX', 'MIN', 'VALUES', 'SET'],
    'bash': ['if', 'then', 'else', 'elif', 'fi', 'for', 'while', 'do', 'done', 'case', 'esac',
             'function', 'return', 'exit', 'echo', 'export', 'local', 'readonly', 'shift',
             'true', 'false', 'in'],
    'yaml': ['true', 'false', 'null', 'yes', 'no', 'on', 'off', 'True', 'False', 'None',
             'YES', 'NO', 'ON', 'OFF', 'NULL', 'Null'],
    'json': ['true', 'false', 'null'],
    'html': ['html', 'head', 'body', 'div', 'span', 'p', 'a', 'img', 'script', 'style',
             'link', 'meta', 'title', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'ul', 'ol', 'li',
             'table', 'tr', 'td', 'th', 'thead', 'tbody', 'form', 'input', 'button', 'textarea',
             'select', 'option', 'label', 'section', 'article', 'header', 'footer', 'nav', 'main',
             'aside', 'iframe', 'canvas', 'svg', 'video', 'audio', 'source', 'br', 'hr'],
    'cpp': ['class', 'struct', 'public', 'private', 'protected', 'virtual', 'const', 'static',
            'void', 'int', 'float', 'double', 'char', 'bool', 'if', 'else', 'for', 'while',
            'do', 'switch', 'case', 'break', 'continue', 'return', 'new', 'delete', 'this',
            'namespace', 'using', 'template', 'typename', 'true', 'false', 'nullptr', 'auto',
            'enum', 'union', 'typedef', 'sizeof', 'const_cast', 'static_cast', 'dynamic_cast',
            'reinterpret_cast', 'try', 'catch', 'throw', 'friend', 'operator', 'inline',
            'explicit', 'mutable', 'extern', 'volatile', 'register', 'signed', 'unsigned',
            'short', 'long', 'constexpr', 'decltype', 'noexcept'],
    'ruby': ['def', 'class', 'module', 'end', 'if', 'elsif', 'else', 'unless', 'while', 'until',
             'for', 'in', 'do', 'case', 'when', 'then', 'break', 'next', 'redo', 'retry', 'return',
             'yield', 'super', 'self', 'nil', 'true', 'false', 'and', 'or', 'not', 'begin', 'rescue',
             'ensure', 'raise', 'attr_reader', 'attr_writer', 'attr_accessor', 'require', 'include',
             'extend', 'alias', 'defined?', 'lambda', 'proc'],
    'php': ['function', 'class', 'interface', 'trait', 'namespace', 'use', 'extends', 'implements',
            'public', 'private', 'protected', 'static', 'final', 'abstract', 'const', 'var', 'new',
            'if', 'else', 'elseif', 'endif', 'switch', 'case', 'break', 'continue', 'default',
            'while', 'endwhile', 'do', 'for', 'endfor', 'foreach', 'endforeach', 'as', 'return',
            'try', 'catch', 'finally', 'throw', 'echo', 'print', 'isset', 'empty', 'unset',
            'array', 'list', 'true', 'false', 'null', 'require', 'include', 'require_once',
            'include_once', 'global', 'clone', 'instanceof', 'yield', 'from'],
    'kotlin': ['fun', 'val', 'var', 'class', 'object', 'interface', 'data', 'sealed', 'enum',
               'abstract', 'open', 'private', 'protected', 'public', 'internal', 'override',
               'if', 'else', 'when', 'for', 'while', 'do', 'break', 'continue', 'return',
               'try', 'catch', 'finally', 'throw', 'import', 'package', 'as', 'in', 'is',
               'null', 'true', 'false', 'this', 'super', 'companion', 'init', 'constructor',
               'by', 'where', 'suspend', 'inline', 'noinline', 'crossinline', 'reified',
               'lateinit', 'inner', 'const', 'operator', 'infix', 'tailrec', 'vararg'],
    'perl': ['sub', 'my', 'our', 'local', 'use', 'require', 'package', 'if', 'elsif', 'else',
             'unless', 'while', 'until', 'for', 'foreach', 'do', 'next', 'last', 'redo',
             'return', 'goto', 'eval', 'die', 'warn', 'undef', 'defined', 'exists', 'delete',
             'shift', 'unshift', 'push', 'pop', 'splice', 'keys', 'values', 'each', 'map',
             'grep', 'sort', 'reverse', 'chomp', 'chop', 'split', 'join', 'print', 'printf',
             'say', 'open', 'close', 'read', 'write', 'BEGIN', 'END', 'true', 'false'],
    'scala': ['def', 'val', 'var', 'class', 'object', 'trait', 'extends', 'with', 'case',
              'sealed', 'abstract', 'private', 'protected', 'override', 'final', 'lazy',
              'if', 'else', 'match', 'for', 'while', 'do', 'yield', 'return', 'try', 'catch',
              'finally', 'throw', 'import', 'package', 'type', 'new', 'this', 'super',
              'true', 'false', 'null', 'None', 'Some', 'Option', 'Either', 'Left', 'Right'],
    'haskell': ['module', 'where', 'import', 'data', 'type', 'newtype', 'class', 'instance',
                'let', 'in', 'if', 'then', 'else', 'case', 'of', 'do', 'return', 'deriving',
                'infixl', 'infixr', 'infix', 'qualified', 'as', 'hiding', 'forall', 'foreign',
                'True', 'False', 'Nothing', 'Just', 'Maybe', 'Either', 'Left', 'Right'],
}

# Add aliases
SYNTAX_KEYWORDS['js'] = SYNTAX_KEYWORDS['javascript']
SYNTAX_KEYWORDS['ts'] = SYNTAX_KEYWORDS['javascript']
SYNTAX_KEYWORDS['typescript'] = SYNTAX_KEYWORDS['javascript']
SYNTAX_KEYWORDS['py'] = SYNTAX_KEYWORDS['python']
SYNTAX_KEYWORDS['sh'] = SYNTAX_KEYWORDS['bash']
SYNTAX_KEYWORDS['shell'] = SYNTAX_KEYWORDS['bash']
SYNTAX_KEYWORDS['zsh'] = SYNTAX_KEYWORDS['bash']
SYNTAX_KEYWORDS['yml'] = SYNTAX_KEYWORDS['yaml']
SYNTAX_KEYWORDS['c++'] = SYNTAX_KEYWORDS['cpp']
SYNTAX_KEYWORDS['cc'] = SYNTAX_KEYWORDS['cpp']
SYNTAX_KEYWORDS['cxx'] = SYNTAX_KEYWORDS['cpp']
SYNTAX_KEYWORDS['htm'] = SYNTAX_KEYWORDS['html']
SYNTAX_KEYWORDS['rb'] = SYNTAX_KEYWORDS['ruby']
SYNTAX_KEYWORDS['pl'] = SYNTAX_KEYWORDS['perl']
SYNTAX_KEYWORDS['pm'] = SYNTAX_KEYWORDS['perl']
SYNTAX_KEYWORDS['kt'] = SYNTAX_KEYWORDS['kotlin']
SYNTAX_KEYWORDS['kts'] = SYNTAX_KEYWORDS['kotlin']
SYNTAX_KEYWORDS['sc'] = SYNTAX_KEYWORDS['scala']
SYNTAX_KEYWORDS['hs'] = SYNTAX_KEYWORDS['haskell']
SYNTAX_KEYWORDS['lhs'] = SYNTAX_KEYWORDS['haskell']

# GFM table separator: matches "|---|" / "|:---|---:|" / etc.
# The trailing group is `*` (not `+`) so single-column tables are valid GFM.
_TABLE_SEPARATOR_RE = re.compile(r'^\s*\|?\s*:?-{3,}:?\s*(\|\s*:?-{3,}:?\s*)*\|?\s*$')

# --- HackMD/Marp <style> block ---------------------------------------------
#
# Only a tiny, hand-picked subset of CSS is honored:
#   - selectors: keyed by the *trailing token* (`h1`, `code`, `blockquote`,
#     `p`, `a`, `table`) or `""` for the bare `.reveal .slides` slide-default.
#   - properties: the SUPPORTED_CSS_PROPERTIES whitelist below.
# Everything else is silently dropped. This keeps the surface small enough
# that a regex parser is more readable than pulling in tinycss2.
SUPPORTED_CSS_PROPERTIES = frozenset({
    'color', 'background-color',
    'font-size', 'font-family', 'font-weight', 'font-style',
    'text-align', 'text-decoration',
})

# Slide elements we know how to target. The renderer reads from this map.
SUPPORTED_CSS_SELECTORS = frozenset({
    '',           # bare `.reveal .slides`  → body default
    'h1', 'h2', 'h3',
    'p',
    'code', 'pre',
    'a',
    'blockquote',
    'table',
    'li',
})

_STYLE_TAG_RE = re.compile(r'<style\b[^>]*>(.*?)</style>', re.IGNORECASE | re.DOTALL)
_CSS_COMMENT_RE = re.compile(r'/\*.*?\*/', re.DOTALL)
_CSS_RULE_RE = re.compile(r'([^{}]+)\{([^{}]*)\}', re.DOTALL)


def _normalize_selector(selector):
    """Reduce a selector like `.reveal .slides h1` to its trailing token `h1`.

    Returns `""` for `.reveal .slides` alone (slide-body default), and the
    raw last whitespace-separated token otherwise. Unknown tokens are passed
    through; the caller filters them against SUPPORTED_CSS_SELECTORS.
    """
    cleaned = selector.strip().rstrip(',').strip()
    if not cleaned:
        return None
    tokens = cleaned.split()
    if not tokens:
        return None
    last = tokens[-1]
    # `.reveal .slides` with no trailing element → body default
    if last == '.slides' and len(tokens) >= 2 and tokens[-2] == '.reveal':
        return ''
    # Strip a leading combinator we don't model (`>`, `+`, `~`)
    if last in ('>', '+', '~'):
        return None
    return last.lower()


def parse_style_block(css_text):
    """Parse a <style> block (or raw CSS) into {selector_token: {prop: value}}.

    The function tolerates:
      - text with or without `<style>...</style>` wrappers
      - CSS comments
      - extra whitespace / multi-line rules
      - unsupported properties (silently dropped)
      - unsupported selectors (silently dropped)

    Returns an empty dict on empty input or when nothing maps. Never raises
    on malformed CSS — the goal is best-effort theming, not validation.
    """
    if not css_text or not css_text.strip():
        return {}

    # Extract content of <style> tags if present; otherwise treat as raw CSS.
    tag_matches = _STYLE_TAG_RE.findall(css_text)
    css = '\n'.join(tag_matches) if tag_matches else css_text

    # Strip comments before splitting into rules.
    css = _CSS_COMMENT_RE.sub('', css)

    result = {}
    for raw_selector, raw_decls in _CSS_RULE_RE.findall(css):
        token = _normalize_selector(raw_selector)
        if token is None or token not in SUPPORTED_CSS_SELECTORS:
            continue
        decls = {}
        for decl in raw_decls.split(';'):
            if ':' not in decl:
                continue
            prop, _, value = decl.partition(':')
            prop = prop.strip().lower()
            value = value.strip()
            if prop in SUPPORTED_CSS_PROPERTIES and value:
                decls[prop] = value
        if decls:
            # Merge if the same trailing token appears in multiple selectors
            # (e.g. `h1 { color: red }` + `.reveal .slides h1 { font-size: 2em }`).
            result.setdefault(token, {}).update(decls)
    return result


def extract_style_block(body):
    """Lift a top-of-deck <style>...</style> block out of `body`.

    Returns (style_text, body_without_style). The style block must appear
    *before* the first markdown heading (`# `, `## `, `### `) so per-slide
    `<style>` tags deeper in the deck are left alone.

    Returns (None, body) when no qualifying block is found.
    """
    if not body:
        return None, body
    # Find first heading offset; the style must precede it to count as
    # deck-level. Heading detection here is line-prefix only — same rule
    # parse_slide uses.
    first_heading = len(body)
    for m in re.finditer(r'(?m)^(#{1,3})\s', body):
        first_heading = m.start()
        break

    search_window = body[:first_heading]
    m = _STYLE_TAG_RE.search(search_window)
    if not m:
        return None, body
    style_text = m.group(0)
    cleaned = body[:m.start()] + body[m.end():]
    cleaned = cleaned.lstrip('\n')
    cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
    return style_text, cleaned

def hex_to_rgb(hex_color):
    """Convert hex color to RGBColor"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


# --- CSS value parsing -----------------------------------------------------

_CSS_NAMED_COLORS = {
    'black': '000000', 'white': 'FFFFFF', 'red': 'FF0000', 'green': '008000',
    'blue': '0000FF', 'yellow': 'FFFF00', 'gray': '808080', 'grey': '808080',
    'orange': 'FFA500', 'purple': '800080', 'pink': 'FFC0CB', 'brown': 'A52A2A',
    'silver': 'C0C0C0', 'maroon': '800000', 'olive': '808000', 'lime': '00FF00',
    'aqua': '00FFFF', 'teal': '008080', 'navy': '000080', 'fuchsia': 'FF00FF',
}


def css_parse_color(value):
    """Parse a CSS color value to a 6-digit hex string (no leading `#`).

    Returns None if the value is unrecognised. Handles:
      - `#RGB`, `#RRGGBB`
      - `rgb(r, g, b)` with integer components
      - the named colors most authors actually use
    """
    if not value:
        return None
    v = value.strip().lower()
    if v.startswith('#'):
        v = v[1:]
        if len(v) == 3:
            return ''.join(c * 2 for c in v).upper()
        if len(v) == 6:
            return v.upper()
        return None
    if v.startswith('rgb('):
        nums = re.findall(r'\d+', v)
        if len(nums) >= 3:
            r, g, b = (max(0, min(255, int(n))) for n in nums[:3])
            return f'{r:02X}{g:02X}{b:02X}'
        return None
    return _CSS_NAMED_COLORS.get(v, None)


def css_parse_font_size(value):
    """Parse a CSS font-size value to a python-pptx Pt() length.

    Supports `px`, `pt`, `em` (treated as 16px base). Returns None if
    the value is unrecognised. We map `px` 1:1 to `pt` — for slide
    output the difference is negligible and px is what reveal.js
    authors write.
    """
    if not value:
        return None
    m = re.match(r'^\s*([0-9]*\.?[0-9]+)\s*(px|pt|em)?\s*$', value, re.IGNORECASE)
    if not m:
        return None
    num = float(m.group(1))
    unit = (m.group(2) or 'px').lower()
    if unit == 'em':
        num *= 16  # 1em ≈ 16px in reveal.js default sizing
    # px == pt for slide rendering (close enough for the values authors use)
    return Pt(num)


def _relative_luminance(rgb_hex):
    """WCAG relative luminance for an `RRGGBB` color (no `#`)."""
    def _channel(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

    r = int(rgb_hex[0:2], 16)
    g = int(rgb_hex[2:4], 16)
    b = int(rgb_hex[4:6], 16)
    return 0.2126 * _channel(r) + 0.7152 * _channel(g) + 0.0722 * _channel(b)


def css_contrast_ratio(fg_hex, bg_hex):
    """WCAG 2.x contrast ratio between two `RRGGBB` colors.

    Returns a float in [1.0, 21.0]. AA large-text threshold is 3.0;
    AA normal-text threshold is 4.5. We use 4.5 because slide titles
    are not reliably "large text" on small screens.
    """
    l1 = _relative_luminance(fg_hex)
    l2 = _relative_luminance(bg_hex)
    lighter, darker = (l1, l2) if l1 > l2 else (l2, l1)
    return (lighter + 0.05) / (darker + 0.05)


# --- style override application -------------------------------------------

_WCAG_AA_NORMAL = 4.5


def _resolve_slide_bg(colors):
    """Return the slide background as `RRGGBB` (no `#`).

    The converter does not currently set an explicit slide background, so
    we treat the default Office theme white (FFFFFF) as the effective
    background unless `colors['slideBg']` is provided (a future extension
    hook used by tests).
    """
    bg = colors.get('slideBg') or 'FFFFFF'
    return bg.lstrip('#').upper()


def _color_passes_contrast(fg_hex, colors):
    """Return True iff `fg_hex` clears WCAG AA on the slide background."""
    bg = _resolve_slide_bg(colors)
    return css_contrast_ratio(fg_hex.lstrip('#').upper(), bg) >= _WCAG_AA_NORMAL


def _apply_run_style(run, decls, colors, *, allow_color=True, dropped_colors=None):
    """Apply a parsed declarations dict to a single run.

    Only properties we know how to map are honored. `allow_color` lets
    the caller suppress color application (e.g. when contrast guard
    fails) while still picking up font-size etc.
    """
    if not decls:
        return
    fs = decls.get('font-size')
    if fs:
        size = css_parse_font_size(fs)
        if size is not None:
            run.font.size = size
    if allow_color:
        color = decls.get('color')
        if color:
            hex_val = css_parse_color(color)
            if hex_val is not None:
                if _color_passes_contrast(hex_val, colors):
                    run.font.color.rgb = hex_to_rgb(hex_val)
                elif dropped_colors is not None:
                    dropped_colors.add(hex_val)
    ff = decls.get('font-family')
    if ff:
        # Use the first family name, stripping quotes (CSS `font-family: "Foo", sans-serif`).
        first = ff.split(',')[0].strip().strip('"').strip("'")
        if first:
            run.font.name = first
    fw = decls.get('font-weight')
    if fw:
        fw_lc = fw.strip().lower()
        if fw_lc in ('bold', 'bolder') or (fw_lc.isdigit() and int(fw_lc) >= 600):
            run.font.bold = True
        elif fw_lc in ('normal', 'lighter') or (fw_lc.isdigit() and int(fw_lc) < 600):
            run.font.bold = False
    fst = decls.get('font-style')
    if fst:
        fst_lc = fst.strip().lower()
        if fst_lc == 'italic':
            run.font.italic = True
        elif fst_lc == 'normal':
            run.font.italic = False
    td = decls.get('text-decoration')
    if td and 'underline' in td.lower():
        run.font.underline = True

def load_config(input_file):
    """Load configuration from JSON or YAML file"""
    input_path = Path(input_file)
    config_files = [
        input_path.parent / 'config.json',
        input_path.parent / 'slides-config.json',
        Path.cwd() / 'config.json',
    ]
    
    if HAS_YAML:
        config_files.extend([
            input_path.parent / 'config.yaml',
            input_path.parent / 'config.yml',
            input_path.parent / 'slides-config.yaml',
            input_path.parent / 'slides-config.yml',
            Path.cwd() / 'config.yaml',
            Path.cwd() / 'config.yml',
        ])
    
    for config_path in config_files:
        if config_path.exists():
            try:
                content = config_path.read_text()
                if config_path.suffix in ['.yaml', '.yml']:
                    config = yaml.safe_load(content)
                else:
                    config = json.loads(content)
                print(f"Loaded config from {config_path}")
                return config
            except Exception as e:
                print(f"Warning: Could not parse {config_path}: {e}")
    return {}

def parse_inline_formatting(text):
    """Parse markdown inline formatting and return segments"""
    segments = []
    # Pattern for: **bold**, _italic_, `code`, [text](url)
    pattern = r'(\*\*[^*]+\*\*|_[^_]+_|`[^`]+`|\[[^\]]+\]\([^)]+\))'
    
    last_end = 0
    for match in re.finditer(pattern, text):
        if match.start() > last_end:
            segments.append({'text': text[last_end:match.start()], 'type': 'text'})
        
        matched = match.group()
        if matched.startswith('**') and matched.endswith('**'):
            segments.append({'text': matched[2:-2], 'type': 'bold'})
        elif matched.startswith('_') and matched.endswith('_'):
            segments.append({'text': matched[1:-1], 'type': 'italic'})
        elif matched.startswith('`') and matched.endswith('`'):
            segments.append({'text': matched[1:-1], 'type': 'code'})
        elif matched.startswith('['):
            link_match = re.match(r'\[([^\]]+)\]\(([^)]+)\)', matched)
            if link_match:
                segments.append({'text': link_match.group(1), 'type': 'link', 'url': link_match.group(2)})
        
        last_end = match.end()
    
    if last_end < len(text):
        segments.append({'text': text[last_end:], 'type': 'text'})
    
    if not segments:
        segments.append({'text': text, 'type': 'text'})
    
    return segments

def highlight_code(code, lang, colors):
    """Apply syntax highlighting to code and return list of (text, color) tuples"""
    if not lang:
        return [{'text': code, 'color': colors['darkText']}]
    
    lang = lang.lower().strip()
    
    # Handle diff specially
    if lang == 'diff':
        result = []
        for line in code.split('\n'):
            if line.startswith('+'):
                result.append({'text': line + '\n', 'color': colors['syntaxDiffAdd']})
            elif line.startswith('-'):
                result.append({'text': line + '\n', 'color': colors['syntaxDiffDel']})
            else:
                result.append({'text': line + '\n', 'color': colors['darkText']})
        # Remove trailing newline from last segment
        if result and result[-1]['text'].endswith('\n'):
            result[-1]['text'] = result[-1]['text'][:-1]
        return result
    
    keywords = SYNTAX_KEYWORDS.get(lang, [])
    if not keywords:
        return [{'text': code, 'color': colors['darkText']}]
    
    result = []
    lines = code.split('\n')
    
    for line_idx, line in enumerate(lines):
        i = 0
        while i < len(line):
            # Check for strings (double or single quotes)
            if line[i] in '"\'':
                quote = line[i]
                j = i + 1
                while j < len(line) and line[j] != quote:
                    if line[j] == '\\' and j + 1 < len(line):
                        j += 2
                    else:
                        j += 1
                if j < len(line):
                    j += 1
                result.append({'text': line[i:j], 'color': colors['syntaxString']})
                i = j
                continue
            
            # Check for comments
            if line[i:i+2] == '//' or line[i] == '#':
                result.append({'text': line[i:], 'color': colors['syntaxComment']})
                break
            
            # Check for numbers
            if line[i].isdigit():
                j = i
                while j < len(line) and (line[j].isdigit() or line[j] in '.xXabcdefABCDEF'):
                    j += 1
                result.append({'text': line[i:j], 'color': colors['syntaxNumber']})
                i = j
                continue
            
            # Check for keywords/identifiers
            if line[i].isalpha() or line[i] == '_':
                j = i
                while j < len(line) and (line[j].isalnum() or line[j] == '_'):
                    j += 1
                word = line[i:j]
                
                # Check if followed by ( for function detection
                next_char = line[j] if j < len(line) else ''
                
                if word in keywords or word.upper() in keywords:
                    result.append({'text': word, 'color': colors['syntaxKeyword']})
                elif next_char == '(':
                    result.append({'text': word, 'color': colors['syntaxFunction']})
                elif word[0].isupper():
                    result.append({'text': word, 'color': colors['syntaxType']})
                else:
                    result.append({'text': word, 'color': colors['darkText']})
                i = j
                continue
            
            # Default: single character
            result.append({'text': line[i], 'color': colors['darkText']})
            i += 1
        
        # Add newline between lines (except last)
        if line_idx < len(lines) - 1:
            result.append({'text': '\n', 'color': colors['darkText']})
    
    # Merge adjacent segments with same color
    merged = []
    for seg in result:
        if merged and merged[-1]['color'] == seg['color']:
            merged[-1]['text'] += seg['text']
        else:
            merged.append(seg)
    
    return merged

class _SlidesWithStyle(list):
    """List subclass that carries parsed CSS overrides alongside the slides.

    Existing callers iterate over it like a list (it *is* a list); the
    renderer reads `.style_overrides` to pick up deck-level theming.
    """
    style_overrides: dict


def parse_markdown(content):
    """Parse HackMD/Marp markdown into slides"""
    # Remove YAML frontmatter if present
    if content.startswith('---'):
        end_idx = content.find('---', 3)
        if end_idx != -1:
            content = content[end_idx + 3:].strip()

    # Lift the top-of-deck <style> block (if any) before slide splitting
    # so its lines don't leak into the first slide's content.
    style_text, content = extract_style_block(content)
    style_overrides = parse_style_block(style_text) if style_text else {}

    # Split by major sections (---)
    sections = re.split(r'\n---\n', content)
    slides = _SlidesWithStyle()
    slides.style_overrides = style_overrides
    current_section = None
    section_idx = 0
    
    for section in sections:
        # Split by sub-slides (----)
        sub_slides = re.split(r'\n----\n', section)
        
        for i, slide_content in enumerate(sub_slides):
            slide = parse_slide(slide_content.strip())
            
            # Skip empty slides
            if not slide['title'] and not slide['content']:
                continue
            
            # Determine if this is a section slide
            if i == 0 and slide['title'] and slide['title'].startswith('# '):
                current_section = slide['title'].lstrip('# ').strip()
                has_bullet_content = any(item['type'] in ['bullet', 'numbered', 'codeblock'] for item in slide['content'])
                slide['is_section'] = not has_bullet_content
                section_idx += 1
            
            slide['section'] = current_section
            slide['section_idx'] = section_idx
            slides.append(slide)
    
    return slides

def _split_table_row(line):
    """Split a GFM table row on unescaped '|' and trim outer pipes."""
    # Temporarily mask escaped pipes so we can split on real ones
    masked = line.replace(r'\|', '\x00')
    parts = masked.split('|')
    # Drop leading/trailing empty cell from outer pipes
    if parts and parts[0].strip() == '':
        parts = parts[1:]
    if parts and parts[-1].strip() == '':
        parts = parts[:-1]
    return [p.replace('\x00', '|').strip() for p in parts]


def parse_slide(content):
    """Parse individual slide content"""
    lines = content.split('\n')
    slide = {
        'title': None,
        'subtitle': None,
        'content': [],
        'notes': None,
        'is_section': False,
    }

    in_note = False
    in_code_block = False
    code_block_content = []
    code_block_lang = ''
    note_content = []

    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        # Check for note marker
        if stripped.lower().startswith('note:'):
            in_note = True
            note_start = stripped[5:].strip()
            if note_start:
                note_content.append(note_start)
            i += 1
            continue

        if in_note:
            note_content.append(line)
            i += 1
            continue

        # Check for code block
        if stripped.startswith('```'):
            if not in_code_block:
                in_code_block = True
                code_block_lang = stripped[3:].strip()
                code_block_content = []
            else:
                in_code_block = False
                slide['content'].append({
                    'type': 'codeblock',
                    'lang': code_block_lang,
                    'content': '\n'.join(code_block_content),
                })
            i += 1
            continue

        if in_code_block:
            code_block_content.append(line)
            i += 1
            continue

        # Parse headers
        if line.startswith('# '):
            slide['title'] = line
            i += 1
            continue
        if line.startswith('## '):
            slide['title'] = line
            i += 1
            continue
        if line.startswith('### '):
            slide['subtitle'] = line[4:].strip()
            i += 1
            continue

        # Parse GFM tables: a header row like "| a | b |" followed by a
        # separator row like "|---|---|" on the next line. Body rows continue
        # until a blank line or a non-`|` line (GFM termination rules).
        if stripped.startswith('|') and i + 1 < len(lines) and _TABLE_SEPARATOR_RE.match(lines[i + 1]):
            header = _split_table_row(stripped)
            rows = []
            j = i + 2
            while j < len(lines):
                cur = lines[j].strip()
                if not cur:
                    break  # blank line ends the table per GFM
                if not cur.startswith('|'):
                    break  # non-table content ends the table
                rows.append(_split_table_row(cur))
                j += 1
            slide['content'].append({'type': 'table', 'header': header, 'rows': rows})
            i = j
            continue

        # Parse bullet list items (- or *)
        bullet_match = re.match(r'^(\s*)[-*]\s+(.+)', line)
        if bullet_match:
            indent = len(bullet_match.group(1)) // 2
            text = bullet_match.group(2)
            # Handle checkboxes
            if text.startswith('[ ] '):
                text = '☐ ' + text[4:]
            elif text.startswith('[x] ') or text.startswith('[X] '):
                text = '☑ ' + text[4:]
            slide['content'].append({'type': 'bullet', 'text': text, 'indent': indent})
            i += 1
            continue

        # Parse numbered list items
        numbered_match = re.match(r'^(\s*)(\d+)\.\s+(.+)', line)
        if numbered_match:
            indent = len(numbered_match.group(1)) // 2
            num = numbered_match.group(2)
            text = numbered_match.group(3)
            slide['content'].append({'type': 'numbered', 'number': num, 'text': text, 'indent': indent})
            i += 1
            continue

        # Plain text (NOT a bullet - no bullet formatting)
        if stripped:
            slide['content'].append({'type': 'text', 'text': stripped, 'indent': 0})
        i += 1

    if note_content:
        slide['notes'] = '\n'.join(note_content)

    return slide

def add_hyperlink(run, url):
    """Add hyperlink to a run"""
    try:
        rId = run.part.relate_to(url, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink', is_external=True)
        hlinkClick = etree.Element(qn('a:hlinkClick'))
        hlinkClick.set(qn('r:id'), rId)
        run._r.get_or_add_rPr().append(hlinkClick)
    except Exception:
        pass  # Hyperlink failed silently

def add_formatted_runs(paragraph, text, colors, fonts):
    """Add formatted text runs to a paragraph based on markdown formatting"""
    segments = parse_inline_formatting(text)
    
    for seg in segments:
        run = paragraph.add_run()
        run.text = seg['text']
        
        if seg['type'] == 'bold':
            run.font.bold = True
        elif seg['type'] == 'italic':
            run.font.italic = True
        elif seg['type'] == 'code':
            run.font.name = fonts['code']
            run.font.color.rgb = hex_to_rgb(colors['accent'])
        elif seg['type'] == 'link':
            run.font.color.rgb = hex_to_rgb(colors['accent'])
            run.font.underline = True
            add_hyperlink(run, seg['url'])

def disable_bullet(paragraph):
    """Disable bullet formatting for a paragraph"""
    pPr = paragraph._p.get_or_add_pPr()
    # Remove any existing buNone, buChar, buAutoNum elements
    for child in list(pPr):
        if child.tag.endswith(('buNone', 'buChar', 'buAutoNum', 'buFont', 'buClr')):
            pPr.remove(child)
    # Add buNone to explicitly disable bullets
    buNone = etree.SubElement(pPr, qn('a:buNone'))

def style_table_cell(cell, text, *, is_header, colors, fonts):
    """Apply visual styling to a single table cell.

    Header: filled accent background + bold white text.
    Body:   default fill, dark text, normal weight.
    Cell text supports the same inline markdown (bold/italic/code/links)
    as the rest of the deck via add_formatted_runs.
    """
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    # p.clear() drops runs/breaks/fields but preserves <a:pPr>, so any
    # alignment inherited from a table style survives.
    p.clear()
    add_formatted_runs(p, text, colors, fonts)

    if is_header:
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(colors['accent'])
        # Force bold + white on every run so accent-colored `code` runs
        # don't vanish on the accent background. font.name stays conditional
        # so inline `code` keeps its monospace font.
        for run in p.runs:
            run.font.bold = True
            run.font.color.rgb = hex_to_rgb(colors['white'])
            if run.font.name is None:
                run.font.name = fonts['body']
    else:
        # Preserve colors/fonts that add_formatted_runs already chose
        # (accent+monospace for `code`, accent for links); only fill defaults.
        for run in p.runs:
            if run.font.color.type is None:
                run.font.color.rgb = hex_to_rgb(colors['darkText'])
            if run.font.name is None:
                run.font.name = fonts['body']


def add_table_to_slide(slide, table_data, top, colors, fonts):
    """Render a parsed table dict onto `slide` starting at vertical offset `top`.

    `top` is an EMU length (e.g. the result of `Inches(1.5)`), not a float.
    Returns the next y-offset (EMU) the caller should continue from.

    Precondition: `table_data['header']` is a non-empty list. The parser
    guarantees this — a header-less table dict is never emitted.
    """
    header = table_data['header']
    rows = table_data.get('rows', [])
    assert header, "table_data['header'] must be non-empty (parser invariant)"

    n_cols = len(header)
    n_rows = 1 + len(rows)
    # Row height minimum; python-pptx auto-fit will expand if a cell is taller.
    height = Inches(0.45 * n_rows)

    gf = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), top, Inches(9), height)
    table = gf.table

    for c, cell_text in enumerate(header):
        style_table_cell(table.cell(0, c), cell_text, is_header=True, colors=colors, fonts=fonts)

    for r, row in enumerate(rows, start=1):
        for c in range(n_cols):
            cell_text = row[c] if c < len(row) else ''
            style_table_cell(table.cell(r, c), cell_text, is_header=False, colors=colors, fonts=fonts)

    return top + height + Inches(0.15)


def add_section_slide(prs, slide_data, colors, fonts, *, style_overrides=None):
    """Add a section/title slide using Title Slide layout (index 0)"""
    layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(layout)
    style_overrides = style_overrides or {}

    # Get title text
    title_text = slide_data['title'].lstrip('# ').strip() if slide_data['title'] else ''

    # Find and populate placeholders with formatting
    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx
        if idx == 0:  # Title placeholder
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            add_formatted_runs(p, title_text, colors, fonts)
            _apply_overrides_to_paragraph(p, 'h1', style_overrides, colors)
        elif idx == 1:  # Subtitle placeholder
            if slide_data['subtitle']:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                add_formatted_runs(p, slide_data['subtitle'], colors, fonts)
                _apply_overrides_to_paragraph(p, 'h3', style_overrides, colors)
            else:
                shape.text = ''

    # Add speaker notes
    if slide_data['notes']:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data['notes']

    return slide


def _apply_overrides_to_paragraph(paragraph, selector, overrides, colors):
    """Apply `overrides[selector]` (and the bare `""` body default) to every run.

    The bare selector's properties (font-size, text-align, etc.) act as
    *defaults* — applied only if the more specific selector didn't set
    them. This matches CSS specificity intuition without modeling full
    cascade rules.
    """
    if not overrides:
        return
    specific = overrides.get(selector, {})
    default = overrides.get('', {})
    # Build an effective declaration set: specific wins per property.
    effective = {**default, **specific}
    if not effective:
        return
    for run in paragraph.runs:
        _apply_run_style(run, effective, colors)


def _apply_body_overrides(paragraph, overrides, colors, fonts):
    """Apply body-default + inline-`code` overrides to a paragraph's runs.

    The bare `""` selector applies to every run (it's the slide-body
    default). The `code` selector applies only to runs that
    add_formatted_runs flagged as inline code by setting `font.name` to
    the code font.
    """
    if not overrides:
        return
    body_default = overrides.get('', {})
    code_decls = overrides.get('code', {})
    for run in paragraph.runs:
        is_code = run.font.name == fonts['code']
        if is_code and code_decls:
            # Inline code: body default acts as fallback, code wins per property.
            _apply_run_style(run, {**body_default, **code_decls}, colors)
        elif body_default:
            _apply_run_style(run, body_default, colors)

def add_content_slide(prs, slide_data, colors, fonts, *, style_overrides=None):
    """Add a content slide using Title and Content layout (index 1)"""
    layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(layout)
    style_overrides = style_overrides or {}

    # Determine which heading selector this slide's title maps to.
    # Markdown `#` → h1 (used by the title slides); `##` → h2; `###` → h3.
    raw_title = slide_data['title'] or ''
    if raw_title.startswith('### '):
        title_selector = 'h3'
    elif raw_title.startswith('## '):
        title_selector = 'h2'
    else:
        title_selector = 'h1'

    # Get title
    title = slide_data['title']
    if title:
        title = re.sub(r'^#+\s*', '', title).strip()
        title = re.sub(r'^\d+\.\s*', '', title)

    # Find placeholders
    title_shape = None
    body_shape = None

    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx
        if idx == 0:
            title_shape = shape
        elif idx == 1:
            body_shape = shape

    # Set title with formatting
    if title_shape and title:
        tf = title_shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        add_formatted_runs(p, title, colors, fonts)
        _apply_overrides_to_paragraph(p, title_selector, style_overrides, colors)
    
    # Set body content
    if body_shape and slide_data['content']:
        tf = body_shape.text_frame
        
        # Code blocks and tables can't live in the body placeholder (it can't
        # hold a GraphicFrame or styled box); both force explicit positioning.
        needs_explicit_layout = any(
            item['type'] in ('codeblock', 'table') for item in slide_data['content']
        )

        if not needs_explicit_layout:
            # Use body placeholder
            first_para = True

            for item in slide_data['content']:
                if first_para:
                    p = tf.paragraphs[0]
                    first_para = False
                else:
                    p = tf.add_paragraph()

                # Set indent level for bullets/numbered only
                if item['type'] in ['bullet', 'numbered']:
                    p.level = item.get('indent', 0)
                else:
                    # Plain text - disable bullet
                    p.level = 0
                    disable_bullet(p)

                # Build text with number prefix for numbered items
                text = item['text']
                if item['type'] == 'numbered':
                    text = f"{item.get('number', '1')}. {text}"

                # Add formatted runs (bold, italic, code, links)
                add_formatted_runs(p, text, colors, fonts)

                # Apply body-default overrides (the bare `""` selector) to
                # every run, and the `code` selector to inline-code runs only.
                _apply_body_overrides(p, style_overrides, colors, fonts)
        else:
            tf.clear()
            p = tf.paragraphs[0]
            p.text = ""
            
            y_pos = Inches(1.5)
            
            for item in slide_data['content']:
                if item['type'] == 'codeblock':
                    code_lines = item['content'].split('\n')
                    code_height = min(len(code_lines) * 0.22 + 0.3, 3.5)
                    
                    # Background rectangle
                    rect = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0.5), y_pos,
                        Inches(9), Inches(code_height)
                    )
                    rect.fill.solid()
                    rect.fill.fore_color.rgb = hex_to_rgb(colors['codeBlock'])
                    rect.line.color.rgb = hex_to_rgb(colors['codeBorder'])
                    
                    # Code text with syntax highlighting
                    code_box = slide.shapes.add_textbox(
                        Inches(0.6), y_pos + Inches(0.1),
                        Inches(8.8), Inches(code_height - 0.2)
                    )
                    code_tf = code_box.text_frame
                    code_tf.word_wrap = True
                    p = code_tf.paragraphs[0]
                    
                    # Apply syntax highlighting
                    highlighted = highlight_code(item['content'], item['lang'], colors)
                    for seg in highlighted:
                        run = p.add_run()
                        run.text = seg['text']
                        run.font.size = Pt(11)
                        run.font.name = fonts['code']
                        run.font.color.rgb = hex_to_rgb(seg['color'])

                    y_pos += Inches(code_height + 0.15)
                elif item['type'] == 'table':
                    y_pos = add_table_to_slide(slide, item, y_pos, colors, fonts)
                else:
                    # Text content
                    text_box = slide.shapes.add_textbox(
                        Inches(0.5), y_pos,
                        Inches(9), Inches(0.5)
                    )
                    text_tf = text_box.text_frame
                    text_tf.word_wrap = True
                    p = text_tf.paragraphs[0]
                    
                    # Build text with bullet/number prefix
                    text = item['text']
                    if item['type'] == 'bullet':
                        text = '• ' + text
                    elif item['type'] == 'numbered':
                        text = f"{item.get('number', '1')}. {text}"
                    
                    # Add formatted runs
                    add_formatted_runs(p, text, colors, fonts)

                    # Apply body-default + code overrides to this paragraph.
                    _apply_body_overrides(p, style_overrides, colors, fonts)

                    # Set font size for all runs that still don't have one
                    for run in p.runs:
                        if run.font.size is None:
                            run.font.size = Pt(15)
                    
                    y_pos += Inches(0.4)
    
    # Add speaker notes
    if slide_data['notes']:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data['notes']
    
    return slide

def add_sections_to_pptx_file(filepath, sections_info):
    """Add section groupings to a saved PPTX file by modifying its XML"""
    import zipfile
    import tempfile
    import shutil
    
    try:
        # Calculate slide IDs for each section
        slide_id = 256  # PowerPoint slide IDs start at 256
        section_data = []
        for sec_idx in sorted(sections_info.keys()):
            sec_info = sections_info[sec_idx]
            if sec_info['name']:
                slide_ids = list(range(slide_id, slide_id + sec_info['count']))
                section_data.append({
                    'name': sec_info['name'],
                    'id': f'{{0000000{sec_idx}-0000-0000-0000-000000000000}}',
                    'slide_ids': slide_ids
                })
            slide_id += sec_info['count']
        
        if not section_data:
            return
        
        # Read the PPTX file
        with zipfile.ZipFile(filepath, 'r') as z_in:
            pres_xml = z_in.read('ppt/presentation.xml')
            root = etree.fromstring(pres_xml)
            
            # Namespaces
            P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            P14_NS = 'http://schemas.microsoft.com/office/powerpoint/2010/main'
            
            # Find or create extLst
            extLst = root.find(f'{{{P_NS}}}extLst')
            if extLst is None:
                extLst = etree.SubElement(root, f'{{{P_NS}}}extLst')
            
            # Create ext element with section URI
            ext = etree.SubElement(extLst, f'{{{P_NS}}}ext')
            ext.set('uri', '{521415D9-36F7-43E2-AB2F-B90AF26B5E84}')
            
            # Create sectionLst with proper namespace prefix
            # Register the p14 namespace
            etree.register_namespace('p14', P14_NS)
            sectionLst = etree.SubElement(ext, f'{{{P14_NS}}}sectionLst')
            
            # Add sections
            for sec in section_data:
                section = etree.SubElement(sectionLst, f'{{{P14_NS}}}section')
                section.set('name', sec['name'])
                section.set('id', sec['id'])
                
                sldIdLst = etree.SubElement(section, f'{{{P14_NS}}}sldIdLst')
                for sid in sec['slide_ids']:
                    sldId = etree.SubElement(sldIdLst, f'{{{P14_NS}}}sldId')
                    sldId.set('id', str(sid))
            
            # Generate modified XML
            modified_xml = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)
            
            # Write to temp file then replace original
            temp_path = filepath + '.tmp'
            with zipfile.ZipFile(temp_path, 'w', zipfile.ZIP_DEFLATED) as z_out:
                for item in z_in.namelist():
                    if item == 'ppt/presentation.xml':
                        z_out.writestr(item, modified_xml)
                    else:
                        z_out.writestr(item, z_in.read(item))
        
        # Replace original with modified
        shutil.move(temp_path, filepath)
        
    except Exception as e:
        print(f"Note: Could not add section markers: {e}")

def main():
    args = sys.argv[1:]
    input_file = args[0] if args else 'slides.md'
    output_file = args[1] if len(args) > 1 else input_file.replace('.md', '.pptx')
    
    if not Path(input_file).exists():
        print(f"Error: Input file '{input_file}' not found")
        print("Usage: python convert.py <input.md> [output.pptx]")
        sys.exit(1)
    
    # Load config
    config = load_config(input_file)
    colors = {**DEFAULT_COLORS, **config.get('colors', {})}
    fonts = {**DEFAULT_FONTS, **config.get('fonts', {})}
    
    # Read and parse markdown
    content = Path(input_file).read_text(encoding='utf-8')
    slides_data = parse_markdown(content)
    style_overrides = getattr(slides_data, 'style_overrides', {})
    if style_overrides:
        print(f"Loaded style overrides: {sorted(style_overrides.keys())}")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # Track sections for grouping
    section_info = {}

    # Add slides
    for slide_data in slides_data:
        if slide_data['is_section']:
            add_section_slide(prs, slide_data, colors, fonts, style_overrides=style_overrides)
        else:
            add_content_slide(prs, slide_data, colors, fonts, style_overrides=style_overrides)
        
        # Track slides per section
        sec_idx = slide_data.get('section_idx', 0)
        if sec_idx not in section_info:
            section_info[sec_idx] = {'name': slide_data.get('section', 'Section'), 'count': 0}
        section_info[sec_idx]['count'] += 1
    
    # Add sections to presentation for collapsible grouping
    # (done after save by modifying the PPTX file)
    
    # Save presentation
    prs.save(output_file)
    
    # Add section markers to the saved file
    add_sections_to_pptx_file(output_file, section_info)
    
    print(f"Created {output_file} with {len(slides_data)} slides")

if __name__ == '__main__':
    main()
